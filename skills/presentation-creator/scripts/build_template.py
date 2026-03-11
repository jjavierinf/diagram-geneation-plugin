#!/usr/bin/env python3
"""Generate the default.pptx template with 9 example layout slides."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Design tokens
TEXT_COLOR = RGBColor(0x2D, 0x2D, 0x2D)
ACCENT_COLOR = RGBColor(0x00, 0x6D, 0xA3)
BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
PLACEHOLDER_COLOR = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)
FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_DIR = Path(__file__).resolve().parent.parent / "assets" / "templates"


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=TEXT_COLOR, alignment=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    tf.paragraphs[0].space_after = Pt(0)
    try:
        tf.auto_size = None
    except Exception:
        pass
    return tb


def _add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def slide_title(prs):
    """Slide 1 - Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    mid_y = Inches(2.5)
    _add_textbox(slide, Inches(1), mid_y, Inches(11.333), Inches(1),
                 "Presentation Title", 40, bold=True, color=TEXT_COLOR,
                 alignment=PP_ALIGN.CENTER)
    # accent line
    _add_rect(slide, Inches(5.5), mid_y + Inches(1.1), Inches(2.333), Pt(4), ACCENT_COLOR)
    _add_textbox(slide, Inches(1), mid_y + Inches(1.4), Inches(11.333), Inches(0.6),
                 "Subtitle or tagline goes here", 20, color=GRAY,
                 alignment=PP_ALIGN.CENTER)


def slide_section_divider(prs):
    """Slide 2 - Section Divider."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, ACCENT_COLOR)
    _add_textbox(slide, Inches(1), Inches(2.8), Inches(11.333), Inches(1.5),
                 "Section Title", 36, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_content(prs):
    """Slide 3 - Content (text left 60%, image placeholder right 40%)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                 "Content Slide Title", 28, bold=True, color=ACCENT_COLOR)
    _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(7), Inches(5),
                 "Body text goes here. This area takes about 60% of the slide width "
                 "and is intended for explanatory content, key points, or narrative text.",
                 18, color=TEXT_COLOR)
    _add_rect(slide, Inches(8.2), Inches(1.5), Inches(4.5), Inches(5), PLACEHOLDER_COLOR)
    _add_textbox(slide, Inches(8.2), Inches(3.5), Inches(4.5), Inches(0.5),
                 "[ Image ]", 14, color=GRAY, alignment=PP_ALIGN.CENTER)


def slide_two_column(prs):
    """Slide 4 - Two Column."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                 "Two Column Layout", 28, bold=True, color=ACCENT_COLOR)
    _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
                 "Left column content. Use this for one perspective, category, or topic.",
                 18, color=TEXT_COLOR)
    _add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5),
                 "Right column content. Use this for the complementary perspective or topic.",
                 18, color=TEXT_COLOR)


def slide_bullet_list(prs):
    """Slide 5 - Bullet List."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                 "Bullet List", 28, bold=True, color=ACCENT_COLOR)
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    bullets = [
        "First key point or takeaway",
        "Second key point with supporting detail",
        "Third key point to reinforce the message",
        "Fourth key point for completeness",
    ]
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = FONT
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(12)
        p.level = 0


def slide_diagram(prs):
    """Slide 6 - Diagram (small title, large placeholder)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.6),
                 "Diagram Title", 22, bold=True, color=ACCENT_COLOR)
    _add_rect(slide, Inches(0.8), Inches(1.1), Inches(11.7), Inches(5.8), PLACEHOLDER_COLOR)
    _add_textbox(slide, Inches(3), Inches(3.5), Inches(7), Inches(0.5),
                 "[ Diagram / Image Placeholder ]", 16, color=GRAY,
                 alignment=PP_ALIGN.CENTER)


def slide_comparison(prs):
    """Slide 7 - Comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
                 "Comparison", 28, bold=True, color=ACCENT_COLOR)
    # Left
    _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6),
                 "Option A", 22, bold=True, color=ACCENT_COLOR)
    _add_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5),
                 "Description of the first option, its pros, cons, and relevant details.",
                 18, color=TEXT_COLOR)
    # Right
    _add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.6),
                 "Option B", 22, bold=True, color=ACCENT_COLOR)
    _add_textbox(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4.5),
                 "Description of the second option, its pros, cons, and relevant details.",
                 18, color=TEXT_COLOR)


def slide_quote(prs):
    """Slide 8 - Quote."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.333), Inches(2),
                 '\u201cThis is an inspiring or insightful quote that supports the narrative.\u201d',
                 28, color=TEXT_COLOR, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10.333), Inches(0.5),
                 "\u2014 Attribution Name, Role", 16, color=GRAY,
                 alignment=PP_ALIGN.CENTER)


def slide_closing(prs):
    """Slide 9 - Closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    mid_y = Inches(2.3)
    _add_textbox(slide, Inches(1), mid_y, Inches(11.333), Inches(1),
                 "Thank You", 40, bold=True, color=TEXT_COLOR,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), mid_y + Inches(1.1), Inches(11.333), Inches(0.6),
                 "We appreciate your time and attention", 20, color=GRAY,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), mid_y + Inches(2), Inches(11.333), Inches(0.5),
                 "contact@example.com", 16, color=ACCENT_COLOR,
                 alignment=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_section_divider(prs)
    slide_content(prs)
    slide_two_column(prs)
    slide_bullet_list(prs)
    slide_diagram(prs)
    slide_comparison(prs)
    slide_quote(prs)
    slide_closing(prs)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUT_DIR / "default.pptx"
    prs.save(str(out_path))
    print(f"Saved {out_path} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
