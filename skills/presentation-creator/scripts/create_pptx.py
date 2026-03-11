#!/usr/bin/env python3
"""Assemble a PPTX presentation from a JSON slide spec and optional template."""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Design tokens
TEXT_COLOR = RGBColor(0x2D, 0x2D, 0x2D)
ACCENT_COLOR = RGBColor(0x00, 0x6D, 0xA3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x88, 0x88, 0x88)
FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DEFAULT_TEMPLATE = Path(__file__).resolve().parent.parent / "assets" / "templates" / "default.pptx"

# Margins
MARGIN_LEFT = Inches(0.8)
MARGIN_TOP = Inches(0.6)
CONTENT_W = Inches(11.733)  # 13.333 - 2*0.8
CONTENT_H = Inches(5.5)


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=TEXT_COLOR, alignment=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    """Add a styled textbox to a slide."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = FONT
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment

    for p in tf.paragraphs:
        p.space_before = Pt(0)
        p.space_after = Pt(4)

    tf.paragraphs[0].space_before = Pt(0)
    # Set vertical anchor
    try:
        tf.auto_size = None
    except Exception:
        pass
    from pptx.oxml.ns import qn
    tb._element.find(".//" + qn("a:bodyPr")).set("anchor", {
        MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"
    }.get(anchor, "t"))
    return tb


def _add_accent_bar(slide, left, top, width, height):
    """Add accent-colored rectangle."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    return shape


def _add_notes(slide, text):
    """Set speaker notes."""
    if text:
        slide.notes_slide.notes_text_frame.text = text


def _get_image_size(image_path):
    """Get image dimensions; fall back to 16:9 if PIL unavailable."""
    try:
        from PIL import Image
        with Image.open(image_path) as img:
            return img.size  # (width, height)
    except Exception:
        return (1920, 1080)


def _add_image_fit(slide, image_path, left, top, max_w, max_h):
    """Add image preserving aspect ratio within bounding box."""
    p = Path(image_path)
    if not p.exists():
        _add_textbox(slide, left, top, max_w, max_h,
                     f"[Image not found: {image_path}]", font_size=14, color=GRAY)
        return
    w_px, h_px = _get_image_size(image_path)
    aspect = w_px / h_px
    box_aspect = max_w / max_h
    if aspect > box_aspect:
        width = max_w
        height = int(max_w / aspect)
    else:
        height = max_h
        width = int(max_h * aspect)
    slide.shapes.add_picture(str(p), left, top, width, height)


def _clear_slides(prs):
    """Remove all slides from the presentation."""
    from pptx.oxml.ns import qn
    while len(prs.slides) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get(qn("r:id"))
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)


def _blank_slide(prs):
    """Add a blank slide using layout index 6."""
    try:
        layout = prs.slide_layouts[6]
    except IndexError:
        layout = prs.slide_layouts[len(prs.slide_layouts) - 1]
    return prs.slides.add_slide(layout)


# ── Layout builders ──────────────────────────────────────────────────────────

def _build_title(prs, spec):
    slide = _blank_slide(prs)
    # Accent bar at top
    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15))
    # Title centered
    _add_textbox(slide, MARGIN_LEFT, Inches(2.2), CONTENT_W, Inches(1.5),
                 spec.get("title", ""), font_size=40, bold=True, color=ACCENT_COLOR,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Subtitle
    if spec.get("subtitle"):
        _add_textbox(slide, MARGIN_LEFT, Inches(3.8), CONTENT_W, Inches(1.0),
                     spec["subtitle"], font_size=24, color=GRAY,
                     alignment=PP_ALIGN.CENTER)
    return slide


def _build_section_divider(prs, spec):
    slide = _blank_slide(prs)
    # Full accent background
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ACCENT_COLOR
    bg.line.fill.background()
    _add_textbox(slide, MARGIN_LEFT, Inches(2.8), CONTENT_W, Inches(1.5),
                 spec.get("title", ""), font_size=36, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def _build_content(prs, spec):
    slide = _blank_slide(prs)
    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08))
    _add_textbox(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_W, Inches(0.8),
                 spec.get("title", ""), font_size=28, bold=True, color=ACCENT_COLOR)
    body_top = Inches(1.6)
    if spec.get("image") and Path(spec["image"]).exists():
        # Body text on left, image on right
        _add_textbox(slide, MARGIN_LEFT, body_top, Inches(5.5), CONTENT_H,
                     spec.get("body", ""), font_size=18)
        _add_image_fit(slide, spec["image"],
                       Inches(7.0), body_top, Inches(5.5), Inches(4.5))
    else:
        _add_textbox(slide, MARGIN_LEFT, body_top, CONTENT_W, CONTENT_H,
                     spec.get("body", ""), font_size=18)
    return slide


def _build_two_column(prs, spec):
    slide = _blank_slide(prs)
    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08))
    _add_textbox(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_W, Inches(0.8),
                 spec.get("title", ""), font_size=28, bold=True, color=ACCENT_COLOR)
    col_w = Inches(5.6)
    body_top = Inches(1.6)
    _add_textbox(slide, MARGIN_LEFT, body_top, col_w, CONTENT_H,
                 spec.get("left", ""), font_size=18)
    _add_textbox(slide, Inches(7.0), body_top, col_w, CONTENT_H,
                 spec.get("right", ""), font_size=18)
    return slide


def _build_bullet_list(prs, spec):
    return _build_content(prs, spec)


def _build_diagram(prs, spec):
    slide = _blank_slide(prs)
    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08))
    _add_textbox(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_W, Inches(0.8),
                 spec.get("title", ""), font_size=28, bold=True, color=ACCENT_COLOR)
    img_top = Inches(1.6)
    if spec.get("image"):
        _add_image_fit(slide, spec["image"],
                       MARGIN_LEFT, img_top, CONTENT_W, Inches(5.2))
    else:
        _add_textbox(slide, MARGIN_LEFT, img_top, CONTENT_W, Inches(5.2),
                     spec.get("body", "[No diagram image provided]"),
                     font_size=18, color=GRAY)
    return slide


def _build_comparison(prs, spec):
    slide = _blank_slide(prs)
    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08))
    _add_textbox(slide, MARGIN_LEFT, MARGIN_TOP, CONTENT_W, Inches(0.8),
                 spec.get("title", ""), font_size=28, bold=True, color=ACCENT_COLOR)
    col_w = Inches(5.6)
    header_top = Inches(1.6)
    body_top = Inches(2.3)
    if spec.get("left_title"):
        _add_textbox(slide, MARGIN_LEFT, header_top, col_w, Inches(0.6),
                     spec["left_title"], font_size=20, bold=True, color=ACCENT_COLOR)
    if spec.get("right_title"):
        _add_textbox(slide, Inches(7.0), header_top, col_w, Inches(0.6),
                     spec["right_title"], font_size=20, bold=True, color=ACCENT_COLOR)
    _add_textbox(slide, MARGIN_LEFT, body_top, col_w, Inches(4.5),
                 spec.get("left", ""), font_size=18)
    _add_textbox(slide, Inches(7.0), body_top, col_w, Inches(4.5),
                 spec.get("right", ""), font_size=18)
    return slide


def _build_quote(prs, spec):
    slide = _blank_slide(prs)
    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08))
    quote_text = f"\u201c{spec.get('quote', '')}\u201d"
    _add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.333), Inches(2.5),
                 quote_text, font_size=28, bold=False, color=TEXT_COLOR,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if spec.get("attribution"):
        _add_textbox(slide, Inches(1.5), Inches(4.8), Inches(10.333), Inches(0.8),
                     f"\u2014 {spec['attribution']}", font_size=18, color=GRAY,
                     alignment=PP_ALIGN.CENTER)
    return slide


def _build_closing(prs, spec):
    slide = _blank_slide(prs)
    _add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15))
    _add_textbox(slide, MARGIN_LEFT, Inches(2.5), CONTENT_W, Inches(1.5),
                 spec.get("title", "Thank You"), font_size=40, bold=True,
                 color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    if spec.get("contact"):
        _add_textbox(slide, MARGIN_LEFT, Inches(4.2), CONTENT_W, Inches(1.0),
                     spec["contact"], font_size=20, color=GRAY,
                     alignment=PP_ALIGN.CENTER)
    return slide


LAYOUT_BUILDERS = {
    "title": _build_title,
    "section-divider": _build_section_divider,
    "content": _build_content,
    "two-column": _build_two_column,
    "bullet-list": _build_bullet_list,
    "diagram": _build_diagram,
    "comparison": _build_comparison,
    "quote": _build_quote,
    "closing": _build_closing,
}


def create_presentation(spec_path, output_path, template_path=None):
    """Build a PPTX from a JSON spec file."""
    template = Path(template_path) if template_path else DEFAULT_TEMPLATE
    if template.exists():
        prs = Presentation(str(template))
        _clear_slides(prs)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

    with open(spec_path, "r") as f:
        spec = json.load(f)

    for slide_spec in spec.get("slides", []):
        layout = slide_spec.get("layout", "content")
        builder = LAYOUT_BUILDERS.get(layout, _build_content)
        slide = builder(prs, slide_spec)
        _add_notes(slide, slide_spec.get("notes"))

    prs.save(str(output_path))
    print(f"Created {output_path} with {len(spec.get('slides', []))} slides")


def main():
    parser = argparse.ArgumentParser(description="Assemble PPTX from JSON spec")
    parser.add_argument("spec", help="Path to JSON slide spec file")
    parser.add_argument("output", help="Output .pptx path")
    parser.add_argument("--template", default=None, help="Path to template .pptx")
    args = parser.parse_args()
    create_presentation(args.spec, args.output, args.template)


if __name__ == "__main__":
    main()
