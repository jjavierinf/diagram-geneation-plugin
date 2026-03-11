#!/usr/bin/env python3
"""
svg2pptx.py - Convert HTML/SVG architecture diagrams to PowerPoint presentations.

Usage:
    python3 svg2pptx.py <html-path> <output-pptx> [--title "Custom Title"] [--validate | --no-validate]

Screenshots each .diagram-wrap SVG at 4x resolution using Puppeteer,
then assembles a PPTX with a title slide + one slide per diagram.

By default, runs Stage 1 SVG validation (validate_svg.py) before screenshotting.
Use --no-validate to skip validation.
"""

import argparse
import json
import os
import subprocess
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOT_JS = os.path.join(SCRIPT_DIR, 'screenshot_svgs.js')
VALIDATE_SVG = os.path.join(SCRIPT_DIR, 'validate_svg.py')

# 16:9 slide dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def run_svg_validation(html_path):
    """Run Stage 1 SVG lint validation. Returns True if OK to proceed."""
    print(f"Running SVG validation on {html_path}...", file=sys.stderr)
    result = subprocess.run(
        [sys.executable, VALIDATE_SVG, html_path],
        capture_output=False,  # let output go to stderr/stdout for user to see
    )
    if result.returncode == 1:
        print("\nSVG validation FAILED. Fix the issues above before exporting to PPTX.",
              file=sys.stderr)
        print("Use --no-validate to skip validation.", file=sys.stderr)
        return False
    elif result.returncode == 2:
        print("\nSVG validation passed with warnings. Proceeding with PPTX export.",
              file=sys.stderr)
        return True
    else:
        print("SVG validation passed.", file=sys.stderr)
        return True


def screenshot_svgs(html_path, output_dir):
    """Run screenshot_svgs.js and return the manifest."""
    result = subprocess.run(
        ['node', SCREENSHOT_JS, html_path, output_dir],
        capture_output=True, text=True,
        cwd=os.path.dirname(os.path.abspath(html_path))
    )
    if result.returncode != 0:
        print(f"screenshot_svgs.js failed:\n{result.stderr}", file=sys.stderr)
        sys.exit(1)
    return json.loads(result.stdout)


def add_title_slide(prs, title, subtitle=None):
    """Add a dark title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)

    # Title text
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
        p2.alignment = PP_ALIGN.CENTER


def add_diagram_slide(prs, image_path, title, img_width, img_height):
    """Add a slide with section header + full-bleed diagram image."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Section title at top
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0x0F, 0x17, 0x2A)
    p.font.bold = True

    # Calculate image placement (fill width, maintain aspect ratio)
    margin_x = Inches(0.3)
    available_width = SLIDE_WIDTH - 2 * margin_x
    available_height = SLIDE_HEIGHT - Inches(1.0)  # Leave room for title

    # Aspect ratio from SVG viewBox
    aspect = img_width / img_height if img_height > 0 else 16 / 9

    # Fit to available space
    if available_width / available_height > aspect:
        # Height-constrained
        final_height = available_height
        final_width = int(final_height * aspect)
    else:
        # Width-constrained
        final_width = available_width
        final_height = int(final_width / aspect)

    # Center horizontally
    left = (SLIDE_WIDTH - final_width) // 2
    top = Inches(0.9)

    slide.shapes.add_picture(image_path, left, top, final_width, final_height)


def main():
    parser = argparse.ArgumentParser(description='Convert HTML/SVG diagrams to PPTX')
    parser.add_argument('html_path', help='Path to the HTML file with SVG diagrams')
    parser.add_argument('output_pptx', help='Output PPTX file path')
    parser.add_argument('--title', help='Override presentation title')
    validate_group = parser.add_mutually_exclusive_group()
    validate_group.add_argument('--validate', action='store_true', default=True,
                                help='Run SVG validation before export (default)')
    validate_group.add_argument('--no-validate', action='store_true',
                                help='Skip SVG validation')
    args = parser.parse_args()

    html_path = os.path.abspath(args.html_path)
    if not os.path.exists(html_path):
        print(f"Error: {html_path} not found", file=sys.stderr)
        sys.exit(1)

    # Stage 1: SVG validation
    if args.no_validate:
        print("Skipping SVG validation (--no-validate).", file=sys.stderr)
    else:
        if not run_svg_validation(html_path):
            sys.exit(1)

    with tempfile.TemporaryDirectory() as tmpdir:
        # Screenshot all SVGs
        print(f"Screenshotting SVGs from {html_path}...", file=sys.stderr)
        manifest = screenshot_svgs(html_path, tmpdir)

        title = args.title or manifest.get('title', 'Architecture Diagrams')

        # Build PPTX
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        add_title_slide(prs, title, "Architecture & Data Flow")

        for d in manifest['diagrams']:
            img_path = os.path.join(tmpdir, d['file'])
            if os.path.exists(img_path):
                add_diagram_slide(
                    prs, img_path, d['title'],
                    d.get('width', 1100), d.get('height', 600)
                )
                print(f"  Added slide: {d['title']}", file=sys.stderr)

        prs.save(args.output_pptx)
        print(f"Saved {args.output_pptx} ({len(prs.slides)} slides)", file=sys.stderr)


if __name__ == '__main__':
    main()
